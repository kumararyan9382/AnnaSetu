import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_sih_presentation():
    template_path = r"C:\Users\kumar\Documents\Aryan\Edit\SIH2026-IDEA-Presentation-Format.pptx"
    output_pptx_dir = r"C:\Users\kumar\OneDrive\Documents\Aryan\Hackathon"
    os.makedirs(output_pptx_dir, exist_ok=True)
    output_pptx = os.path.join(output_pptx_dir, "AnnaSetu_SIH2026_Idea_Presentation.pptx")
    alt_output_pptx = r"C:\Users\kumar\Documents\Aryan\Edit\AnnaSetu_SIH2026_Idea_Presentation.pptx"

    prs = Presentation(template_path)

    # Color definitions
    DARK_GREEN = RGBColor(12, 74, 40)
    FOREST_GREEN = RGBColor(27, 94, 32)
    CHARCOAL = RGBColor(31, 41, 55)
    NAVY = RGBColor(15, 23, 42)

    # ==========================================
    # SLIDE 1: TITLE PAGE
    # ==========================================
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "TITLE PAGE" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "SMART INDIA HACKATHON 2026"
                p.font.bold = True
                p.font.size = Pt(22)
                p.font.color.rgb = DARK_GREEN
            elif "SMART INDIA HACKATHON" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "AnnaSetu (अन्नसेतु)"
                p.font.bold = True
                p.font.size = Pt(26)
                p.font.color.rgb = FOREST_GREEN
            elif "Problem Statement ID" in text:
                tf = shape.text_frame
                tf.clear()
                
                items = [
                    ("Problem Statement ID:", " SIH26032"),
                    ("Problem Statement Title:", " Smart Automation of Agricultural Procurement, Queue Scheduling, Multi-Mandi Price Discovery & Direct DBT Tracking"),
                    ("Theme:", " Smart Automation"),
                    ("PS Category:", " Software"),
                    ("Sponsoring Ministry:", " Ministry of Consumer Affairs, Food & Public Distribution"),
                    ("Team ID:", " [Your Team ID]"),
                    ("Team Name:", " AnnaSetu Innovators")
                ]
                
                for idx, (label, val) in enumerate(items):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    run1 = p.add_run()
                    run1.text = label
                    run1.font.bold = True
                    run1.font.size = Pt(11)
                    run1.font.color.rgb = DARK_GREEN
                    
                    run2 = p.add_run()
                    run2.text = val
                    run2.font.bold = False
                    run2.font.size = Pt(11)
                    run2.font.color.rgb = CHARCOAL
                    p.space_after = Pt(4)

    # ==========================================
    # SLIDE 2: PROPOSED SOLUTION
    # ==========================================
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "IDEA TITLE" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "AnnaSetu – Smart Farmer Procurement & Multi-Mandi Scheduling Platform"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GREEN
            elif "Proposed Solution" in text or "Detailed explanation" in text:
                tf = shape.text_frame
                tf.clear()
                
                sections = [
                    ("1. Proposed Solution & Architecture:", [
                        ("Smart Capacity-Aware Slot Booking:", " Staggers farmer arrivals into 2-hour windows based on mandi weighbridge capacity, completely ending 24-36 hr tractor queues."),
                        ("Dual Role Gateways:", " Dedicated Farmer Portal (Booking, Live Token Tracking & DBT Passbook) + Mandi Staff Operator Board (Gate Entry, Electronic Weighbridge & Lab)."),
                        ("Multi-Mandi Real-Time Spot Price Aggregator:", " Displays live spot prices across Govt APMC and Private Mandis (ITC e-Choupal, Adani Silos, Reliance Kisan).")
                    ]),
                    ("2. How it Addresses the Problem:", [
                        ("Zero Waiting Bottlenecks:", " Replaces chaotic roadside queues with an automated live token queue (#148)."),
                        ("Eliminates Distress Sales:", " Farmers see verified spot rates before leaving their village, preventing 20% distress selling below MSP."),
                        ("100% Transparent Weighing & Payout:", " Direct electronic scale logging with instant Aadhaar-linked PFMS Direct Benefit Transfer (DBT).")
                    ]),
                    ("3. Innovation & Uniqueness:", [
                        ("Free Kisan AI Mandi Advisor:", " Recommends highest-paying mandi by calculating net profit after deducting round-trip diesel transport costs."),
                        ("Full-Duplex WebSockets Live Sync:", " Sub-50ms live updates to farmer's screen when staff approves each stage with spoken Hindi voice audio.")
                    ])
                ]
                
                for s_title, bullet_items in sections:
                    p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
                    r = p.add_run()
                    r.text = s_title
                    r.font.bold = True
                    r.font.size = Pt(10)
                    r.font.color.rgb = FOREST_GREEN
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    
                    for b_label, b_text in bullet_items:
                        bp = tf.add_paragraph()
                        bp.level = 1
                        r1 = bp.add_run()
                        r1.text = "• " + b_label
                        r1.font.bold = True
                        r1.font.size = Pt(8.5)
                        r1.font.color.rgb = CHARCOAL
                        
                        r2 = bp.add_run()
                        r2.text = b_text
                        r2.font.bold = False
                        r2.font.size = Pt(8.5)
                        r2.font.color.rgb = CHARCOAL
                        bp.space_after = Pt(1)

    # ==========================================
    # SLIDE 3: TECHNICAL APPROACH
    # ==========================================
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "TECHNICAL APPROACH" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GREEN
            elif "Technologies to be used" in text or "Methodology" in text:
                tf = shape.text_frame
                tf.clear()
                
                tech_sections = [
                    ("1. Technologies & Tech Stack:", [
                        ("Backend Engine:", " FastAPI (High-performance asynchronous Python event loops, RESTful OpenAPI /docs)."),
                        ("Real-Time Pub/Sub:", " Full-Duplex WebSockets Manager for sub-50ms queue synchronization across all active devices."),
                        ("Database & Security:", " SQLite / PostgreSQL with ACID transactional integrity, masked Aadhaar, and cryptographic token hashing."),
                        ("Frontend UI & Voice:", " Modern TailwindCSS, responsive mobile-first Kisan UI ('सरल मोड'), Web Speech Synthesis API."),
                        ("Offline Access:", " Integrated 5-Language IVR Voice Hotline & SMS Gateway for keypad feature phones.")
                    ]),
                    ("2. 5-Stage End-to-End Procurement Workflow:", [
                        ("Stage 1 - Slot Booking (Web/IVR):", " Farmer books delivery slot ➔ Capacity balancer issues time-stamped token."),
                        ("Stage 2 - Gate Check-in:", " Gate officer scans QR/Token ➔ System validates slot & checks in vehicle (Voice announcement)."),
                        ("Stage 3 - Certified Weighbridge:", " Vehicle drives onto scale ➔ Gross weight captured directly from digital indicator."),
                        ("Stage 4 - Quality Moisture Lab:", " Grain assay sampled ➔ Moisture % recorded ➔ Certified Grade A assigned."),
                        ("Stage 5 - Tare Weighing & DBT Payout:", " Empty vehicle weighed ➔ Net Qtl calculated ➔ Instant PFMS bank payout triggered ➔ Digital e-Voucher generated.")
                    ])
                ]
                
                for s_title, bullet_items in tech_sections:
                    p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
                    r = p.add_run()
                    r.text = s_title
                    r.font.bold = True
                    r.font.size = Pt(10)
                    r.font.color.rgb = FOREST_GREEN
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    
                    for b_label, b_text in bullet_items:
                        bp = tf.add_paragraph()
                        bp.level = 1
                        r1 = bp.add_run()
                        r1.text = "• " + b_label
                        r1.font.bold = True
                        r1.font.size = Pt(8.5)
                        r1.font.color.rgb = CHARCOAL
                        
                        r2 = bp.add_run()
                        r2.text = b_text
                        r2.font.bold = False
                        r2.font.size = Pt(8.5)
                        r2.font.color.rgb = CHARCOAL
                        bp.space_after = Pt(1)

    # ==========================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ==========================================
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "FEASIBILITY AND VIABILITY" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "FEASIBILITY, SCALABILITY & RISK MITIGATION"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GREEN
            elif "Analysis of the feasibility" in text or "Potential challenges" in text:
                tf = shape.text_frame
                tf.clear()
                
                feasibility_sections = [
                    ("1. Feasibility & Scalability Analysis:", [
                        ("Lightweight & Low Bandwidth:", " Bundle size < 2 MB, loads in under 300ms, fully operational on 2G/3G networks."),
                        ("Zero Capital Hardware Burden:", " Mandi staff requires only basic Android smartphones/tablets; no proprietary hardware."),
                        ("Plug-and-Play Interoperability:", " Standard REST APIs enable rapid integration with state portals (e-Kharid, Meri Fasal Mera Byora, e-NAM).")
                    ]),
                    ("2. Challenges & Strategic Mitigations:", [
                        ("Rural Internet Barriers:", " MITIGATION: Offline-first IVR Voice Hotline (/ivr) & SMS token dispatch operable with zero internet."),
                        ("Digital Literacy Gap:", " MITIGATION: 1-Click 'सरल मोड' with large touch targets, multilingual audio readouts & Hindi voice guidance."),
                        ("Queue Manipulation & Bribery:", " MITIGATION: Immutable cryptographic token ledger tied to vehicle number & Aadhaar hash."),
                        ("Weighbridge Downtime:", " MITIGATION: Real-time load balancer automatically diverts incoming farmer bookings to nearest active center.")
                    ])
                ]
                
                for s_title, bullet_items in feasibility_sections:
                    p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
                    r = p.add_run()
                    r.text = s_title
                    r.font.bold = True
                    r.font.size = Pt(10)
                    r.font.color.rgb = FOREST_GREEN
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    
                    for b_label, b_text in bullet_items:
                        bp = tf.add_paragraph()
                        bp.level = 1
                        r1 = bp.add_run()
                        r1.text = "• " + b_label
                        r1.font.bold = True
                        r1.font.size = Pt(8.5)
                        r1.font.color.rgb = CHARCOAL
                        
                        r2 = bp.add_run()
                        r2.text = b_text
                        r2.font.bold = False
                        r2.font.size = Pt(8.5)
                        r2.font.color.rgb = CHARCOAL
                        bp.space_after = Pt(1)

    # ==========================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ==========================================
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "IMPACT AND BENEFITS" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "SOCIO-ECONOMIC IMPACT & NATIONAL BENEFITS"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GREEN
            elif "Potential impact" in text or "Benefits of the solution" in text:
                tf = shape.text_frame
                tf.clear()
                
                impact_sections = [
                    ("1. Target Beneficiaries & Scope:", [
                        ("Audience:", " 150+ Million Indian Farmers, APMC Mandi Staff, Private Corporate Hubs, State Civil Supplies & FCI.")
                    ]),
                    ("2. Direct Economic Impact for Farmers:", [
                        ("+15% to 20% Higher Realized Income:", " Farmers gain access to private mandi bonus rates (+Rs. 85/Qtl on Wheat, +Rs. 170/Qtl on Mustard)."),
                        ("Saves Rs. 500 - Rs. 1,200 per Trip:", " Completely eliminates 24-36 hrs of tractor diesel idling in mandi queues."),
                        ("100% Elimination of Middlemen Cuts:", " Full sale value credited directly to farmer's Aadhaar-linked bank account via DBT.")
                    ]),
                    ("3. Social & Environmental Impact:", [
                        ("Dignified 2-Hour Procurement Turnaround:", " Ends grueling multi-day roadside wait times without food or sanitation."),
                        ("Carbon Footprint Reduction:", " Eliminates millions of liters of diesel exhaust from idling tractor fleets."),
                        ("Post-Harvest Spoilage Prevention:", " Rapid direct intake prevents grain deterioration from unseasonal rains and open-air storage.")
                    ]),
                    ("4. Ministry & Policymaker Value:", [
                        ("Real-Time Analytics:", " Live state-wise procurement velocity, weighbridge utilization, and automated anti-hoarding alerts.")
                    ])
                ]
                
                for s_title, bullet_items in impact_sections:
                    p = tf.paragraphs[0] if tf.paragraphs[0].text == "" else tf.add_paragraph()
                    r = p.add_run()
                    r.text = s_title
                    r.font.bold = True
                    r.font.size = Pt(10)
                    r.font.color.rgb = FOREST_GREEN
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    
                    for b_label, b_text in bullet_items:
                        bp = tf.add_paragraph()
                        bp.level = 1
                        r1 = bp.add_run()
                        r1.text = "• " + b_label
                        r1.font.bold = True
                        r1.font.size = Pt(8.5)
                        r1.font.color.rgb = CHARCOAL
                        
                        r2 = bp.add_run()
                        r2.text = b_text
                        r2.font.bold = False
                        r2.font.size = Pt(8.5)
                        r2.font.color.rgb = CHARCOAL
                        bp.space_after = Pt(1)

    # ==========================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ==========================================
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if "RESEARCH" in text:
                tf = shape.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                p.text = "RESEARCH, POLICY ALIGNMENT & REFERENCES"
                p.font.bold = True
                p.font.size = Pt(16)
                p.font.color.rgb = DARK_GREEN
            elif "Details / Links" in text:
                tf = shape.text_frame
                tf.clear()
                
                ref_items = [
                    ("1. SIH 2026 Problem Statement SIH26032:", " Ministry of Consumer Affairs, Food & Public Distribution (Smart Automation)."),
                    ("2. Commission for Agricultural Costs & Prices (CACP):", " Price Policy Reports for Kharif & Rabi Crops 2026-27 (MSP Rate & FAQ moisture standards)."),
                    ("3. e-NAM (National Agriculture Market):", " Guidelines on Unified Agricultural Marketing & Private Mandi / e-Choupal Interoperability."),
                    ("4. Public Financial Management System (PFMS):", " Standard Operating Procedures for Direct Benefit Transfer (DBT) Direct Bank Credit."),
                    ("5. NITI Aayog & FAO Supply Chain Reports:", " 'Transforming Agricultural Logistics, Post-Harvest Management & Mandi Modernization in India'."),
                    ("6. Digital India Bhashini Initiative:", " Multi-lingual Indic speech synthesis frameworks for rural digital accessibility.")
                ]
                
                for idx, (r_label, r_val) in enumerate(ref_items):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    r1 = p.add_run()
                    r1.text = "• " + r_label
                    r1.font.bold = True
                    r1.font.size = Pt(9.5)
                    r1.font.color.rgb = FOREST_GREEN
                    
                    r2 = p.add_run()
                    r2.text = r_val
                    r2.font.bold = False
                    r2.font.size = Pt(9.5)
                    r2.font.color.rgb = CHARCOAL
                    p.space_after = Pt(4)

    # Update Team Name on all slides
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "Your Team Name" in shape.text_frame.text:
                shape.text_frame.text = "AnnaSetu Innovators"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.color.rgb = DARK_GREEN

    # Remove Slide 7 (Instructions slide as per SIH guidelines)
    if len(prs.slides) > 6:
        # To delete slide in python-pptx:
        rId = prs.slides._sldIdLst[6].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[6]

    # Save to both destinations
    prs.save(output_pptx)
    prs.save(alt_output_pptx)
    print("Presentation saved successfully at:")
    print("1.", output_pptx)
    print("2.", alt_output_pptx)

if __name__ == "__main__":
    create_sih_presentation()
